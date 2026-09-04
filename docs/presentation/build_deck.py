#!/usr/bin/env python3
"""Build the NanoLang developer deck from repository evidence."""

from __future__ import annotations

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
REPO = Path(os.environ.get("NANOLANG_DECK_REPO", str(HERE.parents[1])))
OUT = Path(os.environ.get("NANOLANG_DECK_OUTPUT", str(HERE / "nanolang-developer-overview.pptx")))
ASSETS = HERE / "assets"
W, H = 13.333, 7.5
INK = RGBColor(0x10, 0x13, 0x17)
PANEL = RGBColor(0x23, 0x28, 0x2F)
FOG = RGBColor(0xEE, 0xF1, 0xF3)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
GREEN = RGBColor(0x76, 0xB9, 0x00)
BLUE = RGBColor(0x72, 0xB7, 0xD6)
STEEL = RGBColor(0x65, 0x70, 0x7C)


def box(slide, x, y, w, h, fill, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def text(slide, value, x, y, w, h, size=20, color=INK, bold=False, mono=False, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.05)
    frame.margin_top = frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.text = value
    for paragraph in frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Courier New" if mono else "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return shape


def notes(slide, lines):
    slide.notes_slide.notes_text_frame.text = "\n".join(lines)


def footer(slide, number):
    text(slide, "NANOLANG", 0.45, 7.26, 1.2, 0.15, 8, STEEL, True)
    text(slide, str(number), 12.65, 7.26, 0.25, 0.15, 8, STEEL, True, align=PP_ALIGN.RIGHT)


def title(slide, headline, subtitle, number):
    """Headline, subtitle, footer.

    The headline colour follows the slide's own background: drawn in ink it
    disappears entirely on the dark slides, which is what it had been doing on
    every second slide of the deck. `slide()` records the fill it painted so
    this does not have to be passed at each call site and cannot be forgotten
    at one of them.
    """
    dark = getattr(slide, "_nanolang_dark", False)
    text(slide, headline, 0.65, 0.5, 11.8, 0.7, 28, FOG if dark else INK, True)
    text(slide, subtitle, 0.68, 1.28, 11.6, 0.45, 13, BLUE if dark else STEEL)
    footer(slide, number)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def slide(fill=FOG):
        s = prs.slides.add_slide(blank)
        box(s, 0, 0, W, H, fill)
        # Remembered so title() can choose a headline colour that contrasts.
        s._nanolang_dark = (fill == INK)
        return s

    mascot = ASSETS / "nanolang-mascot.png"

    # 1 — cover
    s = slide(INK)
    s.shapes.add_picture(str(mascot), Inches(8.0), Inches(0.0), width=Inches(5.33), height=Inches(7.5))
    box(s, 0, 0, 9.1, H, INK)
    text(s, "NANOLANG 4.0", 0.7, 0.6, 3.0, 0.3, 13, GREEN, True)
    text(s, "I say what I mean.\nI compile myself.\nI show my evidence.", 0.7, 1.55, 7.0, 2.5, 34, FOG, True)
    # Someone meeting this deck cold needs to know what I am before being told
    # what I prove. The previous subtitle assumed both.
    text(s, "A small language designed to be written by machines and audited by humans.\n"
            "A developer's view of my syntax, compiler, NanoISA, NanoVM, and what my verifier proves.",
         0.75, 4.35, 7.5, 1.2, 16, BLUE)
    text(s, "NanoISA v2 · NanoVM v2 · bytecode that is verified, not merely well-formed", 0.75, 6.55, 7.2, 0.3, 13, ORANGE, True)
    notes(s, ["Authority: docs/PERSONA.md, README.md, docs/RELEASE_4.0.md.",
              "I describe tested behavior. Work I have not done is labelled as such."])

    # 2 — syntax
    s = slide(); title(s, "Machines write code now. The bottleneck is checking it.", "So I refuse ambiguity: one canonical form, explicit boundaries, and evidence required to compile.", 2)
    for i, (head, body_) in enumerate([("PREFIX", "(f x y)"), ("TYPES", "int · float · bool · string"), ("PROOF", "shadow fn { ... }")]):
        x = 0.8 + i * 4.15
        box(s, x, 2.4, 3.5, 2.0, PANEL, True); text(s, head, x + .2, 2.7, 3.1, .3, 14, ORANGE, True); text(s, body_, x + .2, 3.25, 3.1, .7, 20, FOG, True, mono=True)
    notes(s, ["Authority: docs/PERSONA.md and docs/CANONICAL_STYLE.md.",
              "This slide carries the thesis. Everything after it is a mechanism serving this claim,",
              "and a reader who does not accept it here will not care about the module format.",
              "PROOF is not decoration: a function without a shadow test does not compile."])

    # 3 — two paths
    s = slide(INK); title(s, "One source language, two execution paths.", "The C path is my native baseline. NanoISA and NanoVM make the intermediate explicit.", 3)
    for i, (head, body_, color) in enumerate([(".nano", "source + shadow tests", ORANGE), ("nanoc", "generated C", BLUE), ("nano_virt", "NVM bytecode", GREEN), ("nano_vm", "verified execution", FOG)]):
        x = .75 + i * 3.05; box(s, x, 2.6, 2.45, 1.35, PANEL, True); text(s, head, x+.15, 2.82, 2.15, .25, 15, color, True, mono=True); text(s, body_, x+.15, 3.2, 2.15, .45, 12, FOG)
    text(s, "compile → verify → execute", 4.15, 5.2, 5, .35, 21, ORANGE, True, mono=True, align=PP_ALIGN.CENTER)
    notes(s, ["Authority: docs/NANOISA.md and src/nanovirt/main.c."])

    # 4 — readable bytecode
    s = slide(); title(s, "NanoISA is readable bytecode, not a hidden intermediate.", "I keep the serialized module inspectable and use decoded forms only inside execution.", 4)
    box(s, .75, 2.1, 5.7, 3.7, PANEL, True); text(s, ".function main 0 0 0 int 1\n  PUSH_I64 40\n  PUSH_I64 2\n  I64_MUL\n  RET", 1.05, 2.55, 5.1, 2.2, 17, FOG, mono=True)
    box(s, 6.8, 2.1, 5.75, 3.7, INK, True); text(s, "schema → metadata → assembler\n\nThe source of truth is\nspec/nanoisa.yaml", 7.15, 2.6, 5.0, 1.8, 20, GREEN, True)
    notes(s, ["Authority: spec/nanoisa.yaml, scripts/gen_nanoisa_schema.py, src/nanoisa/assembler.c, src/nanoisa/disassembler.c.",
              "Canonical disassembly reassembles to byte-identical bytecode; make test-disasm-roundtrip."])

    # 5 — v2 module format
    s = slide(INK); title(s, "My module format carries the instruction set.", "A .nvm file is a v2 container: versioned, bounded, and refused rather than guessed at.", 5)
    for i, (head, body_, color) in enumerate([
            ("BOUNDED", "every range checked\nby subtraction", GREEN),
            ("TYPED", "signatures deduplicated,\nindices comparable", BLUE),
            ("ACYCLIC", "layouts nest only\ndownward", ORANGE),
            ("EXACT", "constants carry\nexplicit lengths", FOG)]):
        x = .75 + i * 3.05; box(s, x, 2.35, 2.45, 2.15, PANEL, True)
        text(s, head, x+.2, 2.6, 2.05, .3, 13, color, True)
        text(s, body_, x+.2, 3.1, 2.05, 1.2, 14, FOG)
    text(s, "ten sections · a malformed module is refused, never guessed at", 2.0, 5.35, 9.3, .4, 17, ORANGE, True, mono=True, align=PP_ALIGN.CENTER)
    notes(s, ["Authority: src/nanoisa/nvm_format_v2.[ch], src/nanoisa/nvm_v2_*.c, docs/NANOISA.md.",
              "Subtraction form is the point: an offset near the top of the range cannot wrap into it.",
              "An explicit length is why a string holding an embedded zero survives a round trip."])

    # 6 — what the verifier proves
    s = slide(); title(s, "My verifier proves a program before it runs.", "Each property below was a run-time trap. Now it is a rejection.", 6)
    props = [("stack height", "through every basic block, merges must agree"),
             ("operand types", "a known contradiction is refused; unknown never fails"),
             ("return shape", "every exit leaves exactly what the function declares"),
             ("operand depth", "producer declares it, loader confirms it"),
             ("ownership", "retain and release must balance on every path")]
    for i, (head, body_) in enumerate(props):
        y = 1.95 + i * .82; box(s, .9, y, 11.5, .66, PANEL, True)
        text(s, head, 1.2, y+.17, 2.9, .3, 15, GREEN, True)
        text(s, body_, 4.3, y+.17, 7.9, .3, 14, FOG)
    notes(s, ["Authority: src/nanoisa/verifier.c, src/nanoisa/verifier_types.c, tests/nanoisa/test_verifier.c (93 tests).",
              "Every call now carries its signature, so a module is provable before it is linked."])

    # 7 — the failure that mattered
    s = slide(INK); title(s, "What my verifier used to miss.", "It declared stack effects for 32 of its 161 instructions and skipped the rest.", 7)
    box(s, .8, 2.05, 6.1, 3.85, PANEL, True)
    text(s, "PUSH_I64 1\nI64_ADD    ; no declared effect\nPOP\nPOP\nPOP        ; stack held one\nRET", 1.1, 2.45, 5.5, 2.6, 17, FOG, mono=True)
    text(s, "verified ok", 1.1, 5.15, 5.5, .4, 18, ORANGE, True, mono=True)
    box(s, 7.25, 2.05, 5.3, 3.85, INK, True)
    text(s, "Skipping an instruction\nalso skipped its successors,\nso the walk stopped there\nand everything after it\nwent unchecked.\n\nAbsence of evidence was\nindistinguishable from proof.", 7.6, 2.45, 4.6, 3.1, 17, FOG, True)
    notes(s, ["Authority: docs/RELEASE_4.0.md and the fix in src/nanoisa/verifier.c.",
              "Fixing it exposed a second bug in the same walk and then eight latent codegen bugs it had hidden.",
              "An unknown effect is now a hard failure: absence of data must not read as proof."])

    # 8 — hostile input
    s = slide(); title(s, "I treat every module as hostile input.", "The verifier is one defence. The parsers underneath it are the other, and they are fuzzed.", 8)
    surfaces = ["decoder", "loader", "verifier", "assembler", "disassembler", "co-process"]
    for i, name in enumerate(surfaces):
        x = .85 + (i % 3) * 3.95; y = 2.05 + (i // 3) * .82
        box(s, x, y, 3.65, .64, PANEL, True)
        text(s, name, x + .25, y + .18, 3.2, .3, 16, GREEN, True, mono=True)
    box(s, .85, 3.85, 11.65, 2.0, PANEL, True)
    text(s, "size > total - offset", 1.15, 4.15, 5.0, .35, 20, ORANGE, True, mono=True)
    text(s, "never  offset + size > total", 6.6, 4.2, 5.6, .3, 15, STEEL, True, mono=True)
    text(s, "The second form wraps, so an offset near the top of the range passes the check that "
            "exists to stop it. Every range in the v2 decoder is written as the first form. Argument "
            "limits agree across imports, traps, direct FFI and the co-process, so no path is the lenient one.",
         1.15, 4.75, 11.0, .95, 15, FOG)
    text(s, "784 lines of fuzz and malformed-input tests, added in 4.0, across all six", 2.0, 6.15, 9.3, .4, 17, ORANGE, True, mono=True, align=PP_ALIGN.CENTER)
    notes(s, ["Authority: tests/nanoisa/test_fuzz_malformed.c, tests/nanovm/test_cop_fuzz.c, tests/fuzzing/README.md.",
              "784 lines of fuzz and malformed-input tests were added in 4.0 across the six surfaces named here.",
              "Wrapping arithmetic was removed from code-range and section validation, not merely guarded.",
              "This slide is the systemic answer to slide 7: a verifier that is correct is still only one layer."])

    # 9 — dispatch
    s = slide(); title(s, "I dispatch through a label table, and keep a portable fallback.", "One copy of 161 handlers, reached two ways. Adopted on measurement, not principle.", 9)
    box(s, .85, 2.15, 5.7, 3.0, PANEL, True)
    text(s, "COMPUTED GOTO", 1.15, 2.45, 5.1, .3, 14, GREEN, True)
    text(s, "one indirect branch per opcode\n\n-4.2% on my Forth interpreter\nagainst a 1.0-1.6% noise band", 1.15, 2.95, 5.1, 1.8, 16, FOG)
    box(s, 7.0, 2.15, 5.5, 3.0, INK, True)
    text(s, "PORTABLE SWITCH", 7.3, 2.45, 4.9, .3, 14, BLUE, True)
    text(s, "any compiler without\nlabels as values\n\n-DNANO_NO_COMPUTED_GOTO", 7.3, 2.95, 4.9, 1.8, 16, FOG)
    text(s, "152 programs, both builds, identical output", 2.0, 5.55, 9.3, .4, 17, ORANGE, True, mono=True, align=PP_ALIGN.CENTER)
    notes(s, ["Authority: src/nanovm/vm.c, docs/NANOISA_MEASUREMENTS.md, make test-dispatch-equivalence.",
              "VM_CASE and VM_NEXT are the only difference, so the two strategies cannot drift in what an instruction does.",
              "The threaded build has no loop around the handlers, so a stray break is a compile error rather than a silent exit."])

    # 10 — shadow tests
    s = slide(INK); title(s, "Every function carries a shadow test.", "The test is an executable statement about behavior, not a coverage ornament.", 10)
    box(s, .8, 2.0, 5.8, 3.8, PANEL, True); text(s, "fn gcd(a: int, b: int) -> int {\n    ...\n}\n\nshadow gcd {\n    assert (== (gcd 48 18) 6)\n}", 1.1, 2.35, 5.2, 2.9, 16, FOG, mono=True)
    text(s, "2,632 NanoISA tests\n621 NanoVM tests\n93 verifier tests\n63 NanoVirt tests", 7.15, 2.35, 4.7, 2.2, 22, GREEN, True)
    notes(s, ["Authority: CONTRIBUTING.md and the current test suites, counted at the v4.0.0 tag.",
              "make test-verify-all-programs additionally verifies every program in tests/, because compiling is not verifying."])

    # 11 — FFI
    s = slide(); title(s, "FFI is an explicit unsafe boundary.", "Imports carry signatures. The co-process can keep foreign code outside my VM process.", 11)
    box(s, .8, 2.25, 3.0, 2.6, PANEL, True); text(s, "NanoLang", 1.05, 2.65, 2.5, .3, 18, FOG, True); text(s, "typed call", 1.05, 3.35, 2.5, .3, 15, BLUE, True)
    box(s, 5.15, 2.25, 3.0, 2.6, PANEL, True); text(s, "NanoVM", 5.4, 2.65, 2.5, .3, 18, FOG, True); text(s, "typed trap", 5.4, 3.35, 2.5, .3, 15, GREEN, True)
    box(s, 9.5, 2.25, 3.0, 2.6, PANEL, True); text(s, "nano_cop", 9.75, 2.65, 2.5, .3, 18, FOG, True); text(s, "foreign process", 9.75, 3.35, 2.5, .6, 15, ORANGE, True)
    text(s, "isolation costs about 48x per crossing — which is what batching exists to amortize", 1.0, 5.6, 11.3, .4, 15, BLUE, align=PP_ALIGN.CENTER)
    notes(s, ["Authority: docs/EXTERN_FFI.md, src/nanovm/vm_ffi.c, src/nanovm/cop_protocol.c, docs/NANOISA_MEASUREMENTS.md.",
              "The FFI boundary is tested; it is not part of the formally verified NanoCore subset."])

    # 12 — cycles
    s = slide(INK); title(s, "I collect cycles, so my backends agree about leaks.", "Reference counting cannot reclaim a cycle, and a cycle is constructible from ordinary NanoLang.", 12)
    box(s, .8, 2.1, 6.0, 3.5, PANEL, True)
    text(s, "struct Node {\n  children: array<Node>\n}\n\nset kids (array_push kids n)", 1.1, 2.5, 5.4, 2.2, 16, FOG, mono=True)
    box(s, 7.15, 2.1, 5.4, 3.5, INK, True)
    text(s, "array_push mutates in place,\nso the field and the local are\none object -- and it now\ncontains its own owner.\n\nForbidding that shape would\nforbid trees and graphs.", 7.5, 2.5, 4.7, 2.8, 16, FOG, True)
    notes(s, ["Authority: src/nanovm/heap_cycles.c, Bacon and Rajan (PLDI 2001), tests/nanovm/test_vm.c.",
              "The generated-C runtime already collected cycles; the VM did not, so the two disagreed about whether a program leaks.",
              "The test that matters most is the one where a reachable cycle must survive and still be readable."])

    # 13 — measurement
    s = slide(); title(s, "What I measured, and what I declined because of it.", "A number without its spread is not evidence. Every claim below carries a noise band.", 13)
    rows = [("cold startup", "17.4 ms", "60-1800x a single execution", ORANGE),
            ("Forth interpreter", "282.7 us", "IQR 1.6%", GREEN),
            ("computed goto", "-4.2%", "accepted: beats the band", GREEN),
            ("split tag stacks", "unmeasurable", "declined: band is 5-10%", STEEL),
            ("trap stack ranges", "no effect", "declined: crossing is 48x the copy", STEEL)]
    for i, (name, num, note, color) in enumerate(rows):
        y = 2.0 + i * .78; box(s, .9, y, 11.5, .62, PANEL, True)
        text(s, name, 1.2, y+.15, 3.4, .3, 15, FOG, True)
        text(s, num, 4.8, y+.15, 2.4, .3, 15, color, True, mono=True)
        text(s, note, 7.4, y+.15, 4.8, .3, 14, FOG)
    notes(s, ["Authority: docs/NANOISA_MEASUREMENTS.md and scripts/benchmark_nanoisa.sh.",
              "Before 4.0 the suite timed one process per sample, so every workload took about 17 ms whether it retired 78 instructions or 32,082.",
              "It could not have detected an interpreter change of any size, which is why these questions stayed open."])

    # 14 — boundary
    s = slide(INK); title(s, "4.0 shipped. Here is what I have not done.", "107 commits since 3.5, and more new test code than new source. A known defect is not an unknown one.", 14)
    box(s, .8, 2.0, 5.65, 3.9, PANEL, True); text(s, "4.0 SHIPPED", 1.1, 2.35, 4.9, .3, 15, GREEN, True)
    text(s, "v2 module format\nstack and type verification\nreturn, depth, ownership\nfuzzed parsing surfaces\ncycle collection\nhonest measurement", 1.1, 2.9, 4.7, 2.6, 18, FOG, True)
    box(s, 6.9, 2.0, 5.65, 3.9, INK, True); text(s, "NOT DONE", 7.2, 2.35, 4.9, .3, 15, ORANGE, True)
    text(s, "header dependencies (#211)\nmodule signing → 5.0\nLLVM and Wasm as\n  NanoISA translators\nForth 2012 → 4.1", 7.2, 2.9, 4.7, 2.6, 18, FOG, True)
    notes(s, ["Authority: docs/RELEASE_4.0.md, docs/ROADMAP.md, the open issue list, and git log v3.5.0..v4.0.0.",
              "For a reader arriving from 3.5, the one-line delta is that my bytecode used to be well-formed and is now verified.",
              "Phase 12 is 78 of 78. The right column is scope or filed defect, not vagueness."])

    # 15 — closing
    s = slide(INK); s.shapes.add_picture(str(mascot), Inches(8.7), Inches(.8), width=Inches(3.8), height=Inches(5.7))
    text(s, "Start with the code.\nThen run the gates.", .75, 1.7, 7.4, 1.4, 34, FOG, True)
    text(s, "read · change · shadow-test · verify · measure", .8, 4.0, 7.5, .4, 18, ORANGE, True, mono=True)
    text(s, "docs/RELEASE_4.0.md · docs/ROADMAP.md · docs/NANOISA_MEASUREMENTS.md", .8, 5.25, 7.6, .4, 12, BLUE, mono=True)
    notes(s, ["Authority: CONTRIBUTING.md, docs/PERSONA.md, docs/ROADMAP.md."])
    footer(s, 15)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


def main() -> None:
    output = build()
    manifest = Path(os.environ.get("OBJ_DIR", str(REPO / "_build"))) / "nanolang-developer-overview" / "capability-manifest.json"
    manifest.parent.mkdir(parents=True, exist_ok=True)
    manifest.write_text(json.dumps({"schema": "nanolang/developer-overview@1", "slides": len(Presentation(str(output)).slides), "local_artifact": str(output)}, indent=2) + "\n")
    print(f"built {len(Presentation(str(output)).slides)} slides -> {output}")


if __name__ == "__main__":
    main()
