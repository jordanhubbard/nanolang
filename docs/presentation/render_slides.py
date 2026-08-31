#!/usr/bin/env python3
"""Rasterize the overview PPTX with Pillow for visual QA.

LibreOffice is not required. This is a layout inspection aid: it draws fills, strokes,
images, and text from python-pptx shapes. It is not a pixel-identical PowerPoint
renderer. Output lives under ignored _build/.
"""

from __future__ import annotations

import os
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

HERE = Path(__file__).resolve().parent
SOURCE = Path(os.environ.get("NANOLANG_DECK_SOURCE", str(HERE)))
REPO = Path(os.environ.get("NANOLANG_DECK_REPO", str(SOURCE.parents[1])))
PPTX = Path(
    os.environ.get("NANOLANG_DECK_OUTPUT")
    or (SOURCE / "nanolang-developer-overview.pptx")
)
OUT_DIR = Path(os.environ.get("OBJ_DIR") or (REPO / "_build")) / "nanolang-developer-overview"
EMU_PER_PX = 9525
W, H = 1280, 720
COLS = 5


def _px(value) -> int:
    return int(round(int(value) / EMU_PER_PX))


def _rgb(color) -> tuple[int, int, int] | None:
    if color is None:
        return None
    try:
        return (int(color[0]), int(color[1]), int(color[2]))
    except (TypeError, IndexError, ValueError):
        return None


def _fill_rgba(shape) -> tuple[int, int, int, int] | None:
    fill = getattr(shape, "fill", None)
    if fill is None:
        return None
    try:
        ftype = fill.type
    except (ValueError, AttributeError):
        return None
    if ftype is None:
        return None
    try:
        fore = fill.fore_color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    rgb = _rgb(fore)
    if rgb is None:
        return None
    alpha = 255
    try:
        srgb = fill.fore_color._xFill.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
        )
        if srgb is not None:
            node = srgb.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha"
            )
            if node is not None:
                alpha = round(int(node.get("val", "100000")) / 100000 * 255)
    except (AttributeError, TypeError, ValueError):
        pass
    return (*rgb, alpha)


def _line_rgb(shape) -> tuple[int, int, int] | None:
    line = getattr(shape, "line", None)
    if line is None:
        return None
    try:
        return _rgb(line.color.rgb)
    except (AttributeError, TypeError, ValueError):
        return None


def _font(size_pt: float, bold: bool) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    size = max(8, int(round(size_pt * 96 / 72)))
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Courier New Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Courier New.ttf",
        "/Library/Fonts/Arial.ttf",
    ]
    for path in candidates:
        if Path(path).is_file():
            try:
                return ImageFont.truetype(path, size)
            except OSError:
                continue
    return ImageFont.load_default()


def _wrap(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> list[str]:
    lines: list[str] = []
    for paragraph in text.split("\n"):
        if not paragraph:
            lines.append("")
            continue
        words = paragraph.split(" ")
        current = ""
        for word in words:
            trial = word if not current else f"{current} {word}"
            if draw.textlength(trial, font=font) <= max_width or not current:
                current = trial
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines or [""]


def _draw_shape(base: Image.Image, shape) -> None:
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    x, y = _px(shape.left), _px(shape.top)
    w, h = max(1, _px(shape.width)), max(1, _px(shape.height))
    box = [x, y, x + w, y + h]
    fill = _fill_rgba(shape)
    stroke = _line_rgb(shape)
    outline = stroke
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        from io import BytesIO

        picture = Image.open(BytesIO(shape.image.blob)).convert("RGBA")
        picture = picture.resize((w, h), Image.Resampling.LANCZOS)
        overlay.paste(picture, (x, y), picture)
        base.alpha_composite(overlay)
        return
    if fill:
        auto = getattr(shape, "auto_shape_type", None)
        if auto == MSO_SHAPE.OVAL:
            draw.ellipse(box, fill=fill, outline=outline)
        elif auto == MSO_SHAPE.CHEVRON:
            pts = [
                (x, y),
                (x + w - max(6, w // 4), y),
                (x + w, y + h // 2),
                (x + w - max(6, w // 4), y + h),
                (x, y + h),
                (x + max(6, w // 4), y + h // 2),
            ]
            draw.polygon(pts, fill=fill, outline=outline)
        elif auto == MSO_SHAPE.ROUNDED_RECTANGLE:
            radius = max(2, min(w, h) // 8)
            draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline)
        else:
            draw.rectangle(box, fill=fill, outline=outline)
    base.alpha_composite(overlay)
    if shape.has_text_frame:
        _draw_text(base, shape, x, y, w, h)


def _draw_text(base: Image.Image, shape, x: int, y: int, w: int, h: int) -> None:
    draw = ImageDraw.Draw(base)
    cursor_y = y
    for paragraph in shape.text_frame.paragraphs:
        runs = list(paragraph.runs)
        raw = paragraph.text
        if not raw and not runs:
            continue
        size = 14.0
        bold = False
        color = (16, 19, 23, 255)
        if runs:
            font_size = runs[0].font.size
            if font_size is not None:
                size = font_size.pt
            bold = bool(runs[0].font.bold)
            try:
                rgb = _rgb(runs[0].font.color.rgb)
                if rgb:
                    color = (*rgb, 255)
            except (AttributeError, TypeError, ValueError):
                pass
        font = _font(size, bold)
        lines = _wrap(draw, raw, font, max(8, w - 4))
        align = str(getattr(paragraph, "alignment", "") or "")
        for line in lines:
            if cursor_y > y + h:
                break
            width = draw.textlength(line, font=font) if line else 0
            tx = x
            if "CENTER" in align:
                tx = x + max(0, (w - int(width)) // 2)
            elif "RIGHT" in align:
                tx = x + max(0, w - int(width))
            draw.text((tx, cursor_y), line, font=font, fill=color)
            cursor_y += int(size * 96 / 72) + 2


def text_box_overlaps(slide, page: int) -> list[str]:
    """Report pairs of text-bearing shapes whose frames overlap.

    Geometry-escape only checks the slide edge. Painted overflow and colliding
    titles are a separate defect; this is a layout-inspection aid, not an oracle.
    """
    boxes: list[tuple[str, int, int, int, int]] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        raw = (shape.text_frame.text or "").strip()
        if not raw:
            continue
        x, y, w, h = _px(shape.left), _px(shape.top), _px(shape.width), _px(shape.height)
        boxes.append((raw.splitlines()[0][:40], x, y, x + w, y + h))
    hits: list[str] = []
    for i, (a, ax0, ay0, ax1, ay1) in enumerate(boxes):
        for b, bx0, by0, bx1, by1 in boxes[i + 1 :]:
            if ax0 < bx1 - 4 and bx0 < ax1 - 4 and ay0 < by1 - 4 and by0 < ay1 - 4:
                hits.append(f"slide {page}: {a!r} overlaps {b!r}")
    return hits


def render_slide(slide) -> Image.Image:
    canvas = Image.new("RGBA", (W, H), (255, 255, 255, 255))
    for shape in slide.shapes:
        try:
            _draw_shape(canvas, shape)
        except Exception:
            continue
    return canvas.convert("RGB")


def contact_sheet(images: list[Image.Image]) -> Image.Image:
    thumb_w, thumb_h = 240, 135
    rows = (len(images) + COLS - 1) // COLS
    sheet = Image.new("RGB", (COLS * thumb_w, rows * thumb_h), (238, 241, 243))
    for index, image in enumerate(images):
        thumb = image.resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        col, row = index % COLS, index // COLS
        sheet.paste(thumb, (col * thumb_w, row * thumb_h))
    return sheet


def main() -> None:
    if not PPTX.is_file():
        raise SystemExit(f"missing presentation: {PPTX}")
    presentation = Presentation(str(PPTX))
    slides_dir = OUT_DIR / "rendered-slides"
    slides_dir.mkdir(parents=True, exist_ok=True)
    rendered: list[Image.Image] = []
    overlap_hits: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        overlap_hits.extend(text_box_overlaps(slide, index))
        image = render_slide(slide)
        rendered.append(image)
        path = slides_dir / f"slide-{index:02d}.png"
        image.save(path, "PNG")
        print(f"rendered {path}")
    sheet = contact_sheet(rendered)
    sheet_path = OUT_DIR / "contact-sheet.png"
    sheet.save(sheet_path, "PNG")
    print(f"contact sheet {sheet_path}")
    if overlap_hits:
        print(f"{len(overlap_hits)} text-frame overlap(s):")
        for hit in overlap_hits:
            print(f"  {hit}")
    else:
        print("no text-frame overlaps detected")


if __name__ == "__main__":
    main()
